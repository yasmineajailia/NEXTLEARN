"""
content.py — course-file text extraction for the RAG index.

Native Python extraction (pypdf / python-docx / python-pptx) that replaces the
JS mammoth + pdf-parse + jszip path. Given a course-file URL (served by the Node
app), it fetches the bytes, extracts text, and splits it into snippets.
"""

import base64
import hashlib
import hmac
import io
import json
import os
import re
import time

import requests
from pptx import Presentation
from pypdf import PdfReader

import docx  # python-docx

FETCH_TIMEOUT_S = 30
MAX_SNIPPETS = 8
MIN_SNIPPET_LEN = 40
MAX_SNIPPET_LEN = 600


def _b64url(data: bytes) -> str:
    return base64.urlsafe_b64encode(data).rstrip(b"=").decode("ascii")


def _service_token() -> str | None:
    """Mints a short-lived session JWT (same HS256 scheme as issueSession() in
    src/middleware/auth.ts, signed with the same AUTH_SECRET) so this
    server-to-server fetch can pass Node's requireAuth on /api/media/*. That
    route only checks for ANY valid session — it doesn't look the id up
    against a real user or check role — so a synthetic "rag-indexer" identity
    is fine. Without this, every course-file fetch got a 401 that the broad
    except below silently swallowed as an empty snippet list, so newly
    uploaded/edited course content was never actually embedded — only
    structural chunks (names, quiz prompts) were. Returns None if AUTH_SECRET
    isn't configured, degrading to that same pre-existing empty-snippets
    behavior rather than crashing the reindex."""
    secret = os.environ.get("AUTH_SECRET")
    if not secret:
        return None
    now = int(time.time())
    header = _b64url(json.dumps({"alg": "HS256", "typ": "JWT"}, separators=(",", ":")).encode())
    # Least privilege: /api/media/* only calls requireAuth, which accepts ANY
    # valid role, so the lowest one is enough — an "admin" token here would
    # carry backoffice authority this indexer has no use for.
    payload = _b64url(json.dumps(
        {"sub": "rag-indexer", "role": "student", "iat": now, "exp": now + 300},
        separators=(",", ":"),
    ).encode())
    signing_input = f"{header}.{payload}".encode()
    signature = _b64url(hmac.new(secret.encode(), signing_input, hashlib.sha256).digest())
    return f"{header}.{payload}.{signature}"


def _clean(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def extract_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def extract_docx(data: bytes) -> str:
    document = docx.Document(io.BytesIO(data))
    return "\n\n".join(p.text for p in document.paragraphs)


def extract_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                out.append(shape.text_frame.text)
    return "\n\n".join(out)


def extract_bytes(data: bytes, name: str) -> str:
    n = (name or "").lower()
    if n.endswith(".pdf"):
        return extract_pdf(data)
    if n.endswith(".docx"):
        return extract_docx(data)
    if n.endswith(".pptx"):
        return extract_pptx(data)
    return ""


def split_snippets(text: str) -> list:
    parts = re.split(r"\n{2,}|(?<=[.!?])\s{2,}", text or "")
    snippets = []
    for part in parts:
        cleaned = _clean(part)
        if len(cleaned) >= MIN_SNIPPET_LEN:
            snippets.append(cleaned[:MAX_SNIPPET_LEN])
        if len(snippets) >= MAX_SNIPPETS:
            break
    return snippets


def snippets_from_url(url: str, base_url: str = "") -> list:
    if not url:
        return []
    full = url if url.startswith("http") else f"{base_url.rstrip('/')}{url}"
    headers = {}
    token = _service_token()
    if token:
        headers["Cookie"] = f"nl_session={token}"
    try:
        r = requests.get(full, timeout=FETCH_TIMEOUT_S, headers=headers)
        r.raise_for_status()
        return split_snippets(extract_bytes(r.content, url))
    except Exception:  # noqa: BLE001 — a bad/missing file must not fail the whole reindex
        return []
